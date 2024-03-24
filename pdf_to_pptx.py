import os
from pptx import Presentation
from PyPDF2 import PdfReader

def extract_text_from_pdf(pdf_file):
    text = ''
    with open(pdf_file, 'rb') as f:
        pdf_reader = PdfReader(f)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
    return text

def create_pptx_from_text(text, pptx_file):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content 레이아웃 선택
    title = slide.shapes.title
    title.text = os.path.splitext(os.path.basename(pptx_file))[0]

    content = slide.placeholders[1]
    content.text = text

    prs.save(pptx_file)
    print(f"PDF를 PPTX로 변환 완료: {pptx_file}")

if __name__ == "__main__":
    pdf_file = "input.pdf"  # 변환할 PDF 파일
    pptx_file = "output.pptx"  # 생성할 PPTX 파일

    extracted_text = extract_text_from_pdf(pdf_file)
    create_pptx_from_text(extracted_text, pptx_file)
